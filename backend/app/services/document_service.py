from docx import Document
from docx.shared import Pt, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from typing import List
from io import BytesIO
from app.models import Project, Section

class DocumentService:
    """Service for generating Word and PowerPoint documents"""
    
    @staticmethod
    def generate_word_document(project: Project, sections: List[Section]) -> BytesIO:
        """Generate a Word document from project and sections"""
        doc = Document()
        
        # Add title
        title = doc.add_heading(project.name, 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add topic as subtitle
        subtitle = doc.add_paragraph(f"Topic: {project.main_topic}")
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
        subtitle_format = subtitle.runs[0].font
        subtitle_format.size = Pt(14)
        subtitle_format.italic = True
        
        # Add a line break
        doc.add_paragraph()
        
        # Add sections
        for section in sorted(sections, key=lambda x: x.position):
            # Add section heading
            doc.add_heading(section.title, level=1)
            
            # Add section content
            if section.content:
                paragraphs = section.content.split('\n\n')
                for para_text in paragraphs:
                    if para_text.strip():
                        para = doc.add_paragraph(para_text.strip())
                        para_format = para.paragraph_format
                        para_format.space_after = Pt(12)
                        para_format.line_spacing = 1.15
            else:
                doc.add_paragraph("[Content not yet generated]")
            
            # Add spacing between sections
            doc.add_paragraph()
        
        # Save to BytesIO
        file_stream = BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)
        
        return file_stream
    
    @staticmethod
    def generate_powerpoint_presentation(project: Project, sections: List[Section]) -> BytesIO:
        """Generate a PowerPoint presentation from project and sections"""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = project.name
        subtitle.text = f"{project.main_topic}"
        
        # Content Slides
        for section in sorted(sections, key=lambda x: x.position):
            # Use title and content layout
            content_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_slide_layout)
            
            # Add title
            title = slide.shapes.title
            title.text = section.title
            
            # Add content
            content_placeholder = slide.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            
            if section.content:
                # Split content into lines/bullet points
                lines = section.content.split('\n')
                
                for i, line in enumerate(lines):
                    line = line.strip()
                    if line:
                        # Remove bullet point if it exists
                        if line.startswith('•') or line.startswith('-') or line.startswith('*'):
                            line = line[1:].strip()
                        
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        
                        p.text = line
                        p.level = 0
                        p.font.size = PptPt(18)
            else:
                p = text_frame.paragraphs[0]
                p.text = "[Content not yet generated]"
        
        # Save to BytesIO
        file_stream = BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        
        return file_stream
    
    @staticmethod
    def get_filename(project: Project) -> str:
        """Generate filename for document"""
        # Clean project name for filename
        clean_name = "".join(c for c in project.name if c.isalnum() or c in (' ', '-', '_')).strip()
        clean_name = clean_name.replace(' ', '_')
        
        extension = 'docx' if project.document_type == 'docx' else 'pptx'
        return f"{clean_name}.{extension}"